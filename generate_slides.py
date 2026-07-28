from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Define brand colors from shared/stylesheets/colors.css
MIDNIGHT_NAVY = RGBColor(0x00, 0x2e, 0x5d)
TWILIGHT_BLUE = RGBColor(0x03, 0x48, 0x8f)
PHOENIX_GOLD = RGBColor(0xff, 0xc8, 0x45)
DESERT_SAND = RGBColor(0xf7, 0xf4, 0xec)

# Set slide width and height to 16:9 ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Function to set background
def set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Function to add a title slide
def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title slide layout (0)
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, MIDNIGHT_NAVY)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = PHOENIX_GOLD
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = DESERT_SAND
    subtitle.text_frame.paragraphs[0].font.name = 'Arial'
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    return slide

# Function to add a content slide
def add_content_slide(title_text, bullet_points, dark_mode=False):
    slide_layout = prs.slide_layouts[1] # Title and content layout (1)
    slide = prs.slides.add_slide(slide_layout)
    
    bg_color = MIDNIGHT_NAVY if dark_mode else DESERT_SAND
    title_color = PHOENIX_GOLD if dark_mode else TWILIGHT_BLUE
    text_color = DESERT_SAND if dark_mode else MIDNIGHT_NAVY
    
    set_bg(slide, bg_color)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.color.rgb = text_color
        p.font.name = 'Arial'
        p.font.size = Pt(28)
        p.level = 0
        if point.startswith("  - "):
            p.level = 1
            p.text = point.replace("  - ", "")
            p.font.size = Pt(24)
            
    return slide

# Add slides
add_title_slide(
    "The Internet and the World Wide Web",
    "CIS 133: Chapter 1\nLaying the Foundation"
)

add_content_slide(
    "Module Overview & Objectives",
    [
        "Goal: Build a working map of the web's territory.",
        "Key Objectives:",
        "  - Understand how the Internet, Web, browsers, and servers work together.",
        "  - Describe how IP addresses, DNS, and URLs locate pages.",
        "  - Compare web sources for reliability, security, and ethical use.",
        "Expected Deliverables: Skills Lab 1A, Quick Checks."
    ]
)

add_content_slide(
    "The Internet vs. The Web",
    [
        "They are NOT the same thing.",
        "The Internet: The global network and physical infrastructure.",
        "  - The 'road system' (copper wire, fiber optics, cell towers).",
        "  - Started in 1969 as ARPANET.",
        "The World Wide Web (The Web): Linked pages accessed via browsers.",
        "  - One type of 'traffic' on the road system.",
        "  - Created in 1991 by Tim Berners-Lee (HTML and HTTP)."
    ],
    dark_mode=True
)

add_content_slide(
    "The Client-Server Model",
    [
        "Every web interaction has two sides:",
        "Server: A computer that stores websites and sends pages.",
        "  - Stays connected 24/7 waiting for requests.",
        "Client: Any device that asks a server for something.",
        "  - Laptops, phones, smart TVs.",
        "The Rule: The client asks. The server answers."
    ]
)

add_content_slide(
    "Browsers: Your Window to the Web",
    [
        "A browser is a program that does three main jobs:",
        "  1. Sends your request to the right server.",
        "  2. Receives the files the server returns.",
        "  3. Renders HTML & CSS into the page you see.",
        "Browser Cache: Local storage of recently downloaded files to speed up repeat visits.",
        "Browser DevTools: Lets developers look 'behind the page' at the code."
    ],
    dark_mode=True
)

add_content_slide(
    "The Round Trip of a Page Request",
    [
        "1. You type an address or tap a link.",
        "2. The browser looks up which server holds the page (via DNS).",
        "3. The browser sends a request across the Internet.",
        "4. The server finds the requested page and sends it back.",
        "5. The browser renders the returned HTML/CSS into the page.",
        "Usually happens in under a second!"
    ]
)

add_content_slide(
    "Addresses: IPs & Domain Names",
    [
        "IP Address: Numeric label identifying a device (e.g., 91.198.174.192)",
        "  - Built for machines and routers.",
        "Domain Name: Human-readable website name (e.g., phoenixcollege.edu)",
        "  - Built for people.",
        "  - Read right-to-left: Top-Level Domain (.edu), Domain (phoenixcollege), Subdomain (www)."
    ],
    dark_mode=True
)

add_content_slide(
    "DNS: The Translation Service",
    [
        "Domain Name System (DNS) is the Internet's address book.",
        "Translates domain names into IP addresses.",
        "How it works:",
        "  - Browser asks DNS server for a name's IP address.",
        "  - DNS server answers with current IP.",
        "  - Browser connects to that IP."
    ]
)

add_content_slide(
    "Deconstructing a URL",
    [
        "URL (Uniform Resource Locator): The complete address.",
        "Example: https://www.phoenixcollege.edu/degrees",
        "  - Protocol: https:// (How to talk to the server)",
        "  - Subdomain: www",
        "  - Domain: phoenixcollege",
        "  - TLD: .edu",
        "  - Path: /degrees (Which specific page to load)"
    ],
    dark_mode=True
)

add_content_slide(
    "Search Engines: How They Work",
    [
        "Search engines are software systems matching queries to pages.",
        "Three Stages:",
        "  1. Crawling: Bots follow links and read pages.",
        "  2. Indexing: Storing found data in a massive searchable catalog.",
        "  3. Ranking: Sorting matches by predicted usefulness.",
        "SEO (Search Engine Optimization): Helping engines read & rank your pages."
    ]
)

add_content_slide(
    "Evaluating Sources: 5 Questions",
    [
        "1. Author: Is a person or organization identifiable?",
        "2. Host: Who runs the site? (.edu, .gov vs unknown .com)",
        "3. Currency: When was it published or updated?",
        "4. Purpose: Is it informing, selling, or persuading?",
        "5. Evidence: Does the page cite verifiable sources?",
        "Pro-tip: Always cross-reference claims across multiple sources!"
    ],
    dark_mode=True
)

add_content_slide(
    "Ethics & Copyright",
    [
        "Copyright is Automatic: Once created, it's owned.",
        "  - Ability to copy != Right to use.",
        "  - Always ask permission or look for a license.",
        "Creative Commons: Public licenses spelling out allowed reuse (e.g., CC BY 4.0).",
        "Developer Ethics: Would you stand behind the page if your name was on it?"
    ]
)

add_content_slide(
    "Security & Privacy",
    [
        "Security (In Transit):",
        "  - HTTPS + TLS encrypts data so it can't be read in transit.",
        "  - Digital Certificates prove a server is who it claims to be.",
        "  - Never send private info over plain HTTP.",
        "Privacy (At Rest):",
        "  - Cookies store data to remember you (logins, carts) but also track you.",
        "  - As developers: Collect only what you need, protect what you collect."
    ],
    dark_mode=True
)

add_content_slide(
    "Accessibility: The Web for Everyone",
    [
        "The web's founding promise is universal access.",
        "Includes people who see, hear, and move differently.",
        "Screen Readers: Software that reads code aloud or in braille.",
        "  - Images need text descriptions.",
        "  - Structure must be logical.",
        "Accessibility is a legal and ethical requirement, not an afterthought."
    ]
)

add_content_slide(
    "Summary & Next Steps",
    [
        "The Internet (roads) vs The Web (traffic).",
        "Clients request, Servers respond.",
        "DNS translates human names (Domains) to machine numbers (IPs).",
        "Always evaluate sources critically (Author, Host, Currency, Purpose, Evidence).",
        "Build ethically, securely (HTTPS), and accessibly.",
        "Up Next: Setting up VS Code and writing your first HTML in Chapter 2!"
    ],
    dark_mode=True
)

prs.save('/Users/vega/Documents/code/textbooks/cis133/CIS133_Chapter1_Slides.pptx')
print("Successfully generated slides: CIS133_Chapter1_Slides.pptx")
